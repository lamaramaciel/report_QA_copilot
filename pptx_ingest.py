"""Extract slide text and hyperlinks from PPTX or Google Slides extractor CSV."""
from __future__ import annotations

import io
import re
from dataclasses import dataclass, asdict
from typing import Iterable, Any

import pandas as pd
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

URL_RE = re.compile(r'https?://[^\s|,\]>)"\']+')
REF_RE = re.compile(r"\[Ref\d*\]", re.IGNORECASE)


@dataclass
class ClaimRecord:
    claim_id: str
    slide_number: int
    slide_title: str
    element_id: str
    element_name: str
    element_type: str
    paragraph_index: int
    raw_text: str
    clean_claim: str
    ref_markers: str
    urls: str
    url_count: int
    mapping_note: str
    surrounding_text: str


def _unique(items: Iterable[str]) -> list[str]:
    out: list[str] = []
    for item in items:
        value = str(item or "").strip().rstrip(".,;)")
        if value and value not in out:
            out.append(value)
    return out


def _visible_urls(text: str) -> list[str]:
    return _unique(URL_RE.findall(text or ""))


def _shape_type_name(shape: Any) -> str:
    try:
        return str(shape.shape_type).split(".")[-1]
    except Exception:
        return type(shape).__name__


def _shape_click_url(shape: Any) -> str:
    try:
        return str(shape.click_action.hyperlink.address or "").strip()
    except Exception:
        return ""


def _run_url(run: Any) -> str:
    try:
        return str(run.hyperlink.address or "").strip()
    except Exception:
        return ""


def _slide_title(slide: Any) -> str:
    try:
        if slide.shapes.title is not None:
            return str(slide.shapes.title.text or "").strip()
    except Exception:
        pass
    for shape in slide.shapes:
        if getattr(shape, "has_text_frame", False):
            text = str(shape.text or "").strip()
            if text:
                return text.splitlines()[0][:180]
    return ""


def _paragraph_records(
    *,
    paragraph: Any,
    slide_number: int,
    slide_title: str,
    element_id: str,
    element_name: str,
    element_type: str,
    paragraph_index: int,
    surrounding_text: str,
    shape_url: str = "",
) -> list[ClaimRecord]:
    raw_text = str(paragraph.text or "").strip()
    if not raw_text:
        return []

    run_links: list[tuple[str, str]] = []
    for run in paragraph.runs:
        url = _run_url(run)
        if url:
            run_links.append((str(run.text or "").strip(), url))

    urls = _unique([u for _, u in run_links] + _visible_urls(raw_text) + ([shape_url] if shape_url else []))
    if not urls:
        return []

    markers = [m.group(0) for m in REF_RE.finditer(raw_text)]
    clean_claim = REF_RE.sub("", raw_text)
    clean_claim = re.sub(r"https?://\S+", "", clean_claim)
    clean_claim = re.sub(r"\s+", " ", clean_claim).strip(" ;:,-")

    linked_marker_count = sum(1 for run_text, _ in run_links if REF_RE.search(run_text or ""))
    if markers and len(urls) == len(markers):
        note = "Direct marker-to-link mapping available"
    elif markers and len(urls) < len(markers):
        note = f"Incomplete mapping: {len(markers)} marker(s), {len(urls)} URL(s)"
    elif markers and len(urls) > len(markers):
        note = f"Extra URLs: {len(markers)} marker(s), {len(urls)} URL(s)"
    elif linked_marker_count:
        note = "Links detected on [Ref] text runs"
    elif shape_url:
        note = "Shape-level hyperlink; claim mapping should be reviewed"
    else:
        note = "Hyperlinks detected without explicit [Ref] markers"

    if not clean_claim or REF_RE.fullmatch(raw_text.replace(" ", "")):
        note = "Reference-only paragraph; claim mapping unclear"

    claim_id = f"S{slide_number}-E{element_id}-P{paragraph_index + 1}"
    return [ClaimRecord(
        claim_id=claim_id,
        slide_number=slide_number,
        slide_title=slide_title,
        element_id=element_id,
        element_name=element_name,
        element_type=element_type,
        paragraph_index=paragraph_index + 1,
        raw_text=raw_text,
        clean_claim=clean_claim,
        ref_markers=" | ".join(markers),
        urls=" | ".join(urls),
        url_count=len(urls),
        mapping_note=note,
        surrounding_text=surrounding_text[:4000],
    )]


def _walk_shapes(shapes: Any) -> Iterable[Any]:
    for shape in shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from _walk_shapes(shape.shapes)
        else:
            yield shape


def extract_from_pptx(data: bytes) -> pd.DataFrame:
    prs = Presentation(io.BytesIO(data))
    rows: list[ClaimRecord] = []

    for s_idx, slide in enumerate(prs.slides, start=1):
        title = _slide_title(slide)
        for shape in _walk_shapes(slide.shapes):
            element_id = str(getattr(shape, "shape_id", ""))
            element_name = str(getattr(shape, "name", ""))
            element_type = _shape_type_name(shape)
            shape_url = _shape_click_url(shape)

            if getattr(shape, "has_text_frame", False):
                surrounding = str(shape.text or "").strip()
                for p_idx, paragraph in enumerate(shape.text_frame.paragraphs):
                    rows.extend(_paragraph_records(
                        paragraph=paragraph,
                        slide_number=s_idx,
                        slide_title=title,
                        element_id=element_id,
                        element_name=element_name,
                        element_type=element_type,
                        paragraph_index=p_idx,
                        surrounding_text=surrounding,
                        shape_url=shape_url,
                    ))

            if getattr(shape, "has_table", False):
                table = shape.table
                header_values = [str(cell.text or "").strip() for cell in table.rows[0].cells] if len(table.rows) else []
                for r_idx, row in enumerate(table.rows, start=1):
                    row_values = [str(cell.text or "").strip() for cell in row.cells]
                    row_context = " | ".join(value for value in row_values if value)
                    for c_idx, cell in enumerate(row.cells, start=1):
                        column_header = header_values[c_idx - 1] if c_idx - 1 < len(header_values) else ""
                        context_parts = []
                        if column_header:
                            context_parts.append(f"Table column: {column_header}")
                        if row_context:
                            context_parts.append(f"Table row: {row_context}")
                        surrounding = "\n".join(context_parts) or str(cell.text or "").strip()
                        for p_idx, paragraph in enumerate(cell.text_frame.paragraphs):
                            rows.extend(_paragraph_records(
                                paragraph=paragraph,
                                slide_number=s_idx,
                                slide_title=title,
                                element_id=f"{element_id}-R{r_idx}C{c_idx}",
                                element_name=f"{element_name} [R{r_idx}C{c_idx}]",
                                element_type="TABLE_CELL",
                                paragraph_index=p_idx,
                                surrounding_text=surrounding,
                                shape_url=shape_url,
                            ))

    columns = [f.name for f in ClaimRecord.__dataclass_fields__.values()]
    return pd.DataFrame([asdict(r) for r in rows], columns=columns)


def extract_from_google_slides_csv(data: bytes) -> pd.DataFrame:
    df = pd.read_csv(io.BytesIO(data), dtype=str).fillna("")
    expected = {
        "Slide", "Slide Title", "Element ID", "Element Name", "Element Type",
        "Paragraph", "Raw Text", "Clean Claim", "Ref Markers", "URLs",
        "Mapping Note", "Surrounding Text",
    }
    missing = expected - set(df.columns)
    if missing:
        raise ValueError("This CSV is not a Slide QA extraction. Missing: " + ", ".join(sorted(missing)))

    out = pd.DataFrame({
        "claim_id": [f"S{s}-E{e}-P{p}" for s, e, p in zip(df["Slide"], df["Element ID"], df["Paragraph"])],
        "slide_number": pd.to_numeric(df["Slide"], errors="coerce").fillna(0).astype(int),
        "slide_title": df["Slide Title"],
        "element_id": df["Element ID"],
        "element_name": df["Element Name"],
        "element_type": df["Element Type"],
        "paragraph_index": pd.to_numeric(df["Paragraph"], errors="coerce").fillna(0).astype(int),
        "raw_text": df["Raw Text"],
        "clean_claim": df["Clean Claim"],
        "ref_markers": df["Ref Markers"],
        "urls": df["URLs"],
        "url_count": df["URLs"].apply(lambda x: len(_unique(str(x).split(" | ")))),
        "mapping_note": df["Mapping Note"],
        "surrounding_text": df["Surrounding Text"],
    })
    return out


def extract_deck(data: bytes, filename: str) -> pd.DataFrame:
    low = filename.lower()
    if low.endswith(".pptx"):
        return extract_from_pptx(data)
    if low.endswith(".csv"):
        return extract_from_google_slides_csv(data)
    raise ValueError("Supported formats: .pptx or the CSV generated by extract_slide_references.gs")
